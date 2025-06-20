import os
import json
import uuid
import logging
import traceback
import asyncio
from datetime import datetime
from typing import Dict, List, Optional, Any
import tempfile

# ロガーの初期化
logger = logging.getLogger(__name__)

# OpenAIのインポート
try:
    from openai import OpenAI
    logger.info(f"OpenAI library imported successfully")
except ImportError:
    logger.error("OpenAIライブラリがインストールされていません")
    raise
except Exception as e:
    logger.error(f"OpenAIライブラリのインポート中にエラーが発生しました: {str(e)}")
    raise

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt

from app.schemas.presentation import Presentation, Slide, PresentationOptions

# OpenAIクライアントの初期化
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    logger.error("OpenAI API Keyが設定されていません。環境変数OPENAI_API_KEYを設定してください。")
    # デバッグ用にすべての環境変数をログに出力（本番環境では削除すること）
    logger.debug(f"利用可能な環境変数: {list(os.environ.keys())}")

# OpenAIクライアントのインスタンス作成
try:
    client = OpenAI(api_key=api_key) if api_key else None
    logger.info(f"OpenAI client initialized: {'設定済み' if client else '未設定'}")
except Exception as e:
    logger.error(f"OpenAIクライアント初期化エラー: {str(e)}")
    client = None

# 一時ファイル保存用のディレクトリ
TEMP_DIR = tempfile.gettempdir()
logger.info(f"一時ファイル保存ディレクトリ: {TEMP_DIR}")

# メモリ内キャッシュ
# 実際のアプリケーションではデータベース等の永続化層を使用する
presentations_cache: Dict[str, Presentation] = {}


async def generate_presentation_from_text(
    text: str, 
    options: Optional[PresentationOptions] = None
) -> Presentation:
    """テキストからプレゼンテーションを生成する"""
    if options is None:
        options = PresentationOptions()
    
    logger.info(f"プレゼンテーション生成開始: テーマ={options.theme}, スライド数={options.slide_count}")
    logger.debug(f"入力テキスト(先頭100文字): {text[:100]}...")
    
    try:
        # GPT-4を使用してテキストからスライド構造を生成
        logger.info("OpenAI APIリクエスト送信中...")
        
        system_prompt = f"""あなたはプレゼンテーションのスペシャリストです。テキストから高品質なプレゼンテーションを作成します。
以下のテーマでスライドを作成してください: {options.theme}
スライド数の目安は約{options.slide_count}枚です。
各スライドにはタイトルとコンテンツのリストを含めてください。
結果はJSON形式で返してください。次の形式に従ってください:
{{
  "slides": [
    {{
      "title": "スライドタイトル",
      "content": ["コンテンツ項目1", "コンテンツ項目2"]
    }}
  ]
}}
"""
        logger.debug(f"System prompt: {system_prompt}")
        
        # OpenAI APIの呼び出し
        try:
            if not client:
                raise ValueError("OpenAI クライアントが初期化されていません")
            
            # 非同期コンテキストでの同期API呼び出し
            response = await asyncio.to_thread(
                client.chat.completions.create,
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": system_prompt
                    },
                    {
                        "role": "user",
                        "content": text
                    }
                ],
                response_format={"type": "json_object"}  # JSON形式を強制
            )
            logger.info("OpenAI APIからレスポンスを受信しました")
            logger.debug(f"Response: {response}")
        except Exception as e:
            logger.error(f"OpenAI API呼び出しエラー: {str(e)}")
            raise ValueError(f"OpenAI APIの呼び出しに失敗しました: {str(e)}")

        # APIレスポンスからスライドデータを取得
        response_content = response.choices[0].message.content
        if not response_content:
            logger.error("OpenAI APIから空のレスポンスを受信")
            raise ValueError("レスポンスの生成に失敗しました")

        logger.debug(f"OpenAI APIレスポンス: {response_content[:300]}...")

        # JSONパース
        try:
            parsed_response = json.loads(response_content)
            slides_data = parsed_response.get("slides", [])
            logger.info(f"パース成功: スライド数={len(slides_data)}")
        except json.JSONDecodeError as e:
            logger.error(f"JSONパースエラー: {str(e)}")
            logger.debug(f"パースできなかったJSON: {response_content}")
            raise ValueError(f"JSONのパースに失敗しました: {str(e)}")
        
        # Slideオブジェクトに変換
        slides = []
        for i, slide_data in enumerate(slides_data):
            try:
                slide = Slide(
                    title=slide_data["title"],
                    content=slide_data["content"],
                    image_url=slide_data.get("image_url")
                )
                slides.append(slide)
            except KeyError as e:
                logger.warning(f"スライド{i}のデータに必須キーがありません: {str(e)}")
                # エラーをスキップして続行
                continue

        # 画像生成と追加（実装予定）
        if options.include_images and api_key:
            logger.info("画像生成機能は現在実装中...")
            # TODO: DALL-Eや他の画像生成APIの実装

        # プレゼンテーションデータの作成
        presentation_id = str(uuid.uuid4())
        presentation = Presentation(
            id=presentation_id,
            slides=slides,
            theme=options.theme,
            created_at=datetime.now().isoformat(),
            download_url=f"/api/presentations/{presentation_id}/download"
        )

        # キャッシュに保存
        presentations_cache[presentation_id] = presentation
        logger.info(f"プレゼンテーション生成完了: ID={presentation_id}, スライド数={len(slides)}")

        return presentation

    except Exception as e:
        logger.error(f"プレゼンテーション生成エラー: {str(e)}")
        logger.debug(f"詳細なエラー: {traceback.format_exc()}")
        raise ValueError(f"プレゼンテーションの生成中にエラーが発生しました: {str(e)}")


def get_presentation_by_id(presentation_id: str) -> Optional[Presentation]:
    """プレゼンテーションIDからプレゼンテーションデータを取得する"""
    presentation = presentations_cache.get(presentation_id)
    if presentation:
        logger.info(f"キャッシュからプレゼンテーションを取得: ID={presentation_id}")
    else:
        logger.warning(f"プレゼンテーションが見つかりません: ID={presentation_id}")
        
        # デモ用にダミーデータを返す（本番ではこの処理を削除）
        if presentation_id:
            dummy_presentation = Presentation(
                id=presentation_id,
                slides=[
                    Slide(
                        title="デモプレゼンテーション",
                        content=["このプレゼンテーションはデモデータです", "実際のデータはキャッシュに保存されていませんでした"]
                    ),
                    Slide(
                        title="テスト用スライド",
                        content=["テスト用の内容項目1", "テスト用の内容項目2", "テスト用の内容項目3"]
                    ),
                    Slide(
                        title="まとめ",
                        content=["実運用時にはデータベースに保存することをお勧めします", "セッションデータの永続化を検討してください"]
                    )
                ],
                theme="modern",
                created_at=datetime.now().isoformat(),
                download_url=f"/api/presentations/{presentation_id}/download"
            )
            # ダミーデータをキャッシュに保存
            presentations_cache[presentation_id] = dummy_presentation
            logger.info(f"デモ用のダミーデータを生成しました: ID={presentation_id}")
            return dummy_presentation
    
    return presentation


async def generate_powerpoint_file(presentation: Presentation) -> str:
    """プレゼンテーションをPowerPointファイルとして生成する"""
    logger.info(f"PowerPointファイル生成開始: ID={presentation.id}, スライド数={len(presentation.slides)}")
    
    try:
        # PowerPointファイルを作成
        pptx = PPTXPresentation()
        
        # テーマに基づいた設定
        theme_settings = get_theme_settings(presentation.theme)
        logger.debug(f"テーマ設定: {theme_settings}")
        
        # 各スライドの生成
        for i, slide_data in enumerate(presentation.slides):
            logger.debug(f"スライド{i+1}作成中: タイトル={slide_data.title}")
            
            # スライドの追加
            try:
                # 最初のスライドはタイトルスライド、それ以外は通常のスライド
                if i == 0:
                    slide_layout = pptx.slide_layouts[0]  # タイトルスライド
                else:
                    slide_layout = pptx.slide_layouts[1]  # タイトルと内容のレイアウト
                
                slide = pptx.slides.add_slide(slide_layout)
                
                # タイトルの設定
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_data.title
                
                # コンテンツの設定（最初のスライド以外）
                if i > 0:
                    try:
                        content_shape = slide.placeholders[1]  # コンテンツのプレースホルダー
                        text_frame = content_shape.text_frame
                        
                        # 各コンテンツ項目を一行ずつ追加
                        for j, content_item in enumerate(slide_data.content):
                            if j == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = content_item
                            p.level = 0  # 箇条書きレベル
                    except Exception as content_error:
                        logger.error(f"コンテンツ設定エラー: {str(content_error)}")
                # 最初のスライドのサブタイトル設定
                elif i == 0 and len(slide_data.content) > 0:
                    try:
                        # サブタイトル用のプレースホルダーを探す
                        for shape in slide.placeholders:
                            if shape.placeholder_format.type == 2:  # サブタイトルのプレースホルダータイプ
                                shape.text = "\n".join(slide_data.content[:2])
                                break
                    except Exception as subtitle_error:
                        logger.error(f"サブタイトル設定エラー: {str(subtitle_error)}")
            except Exception as e:
                logger.error(f"スライド{i+1}の生成中にエラーが発生: {str(e)}")
                # このスライドはスキップして続行
        
        # 一時ファイルに保存
        file_name = f"presentation_{presentation.id}.pptx"
        file_path = os.path.join(TEMP_DIR, file_name)
        logger.debug(f"PowerPointファイルを保存: {file_path}")
        pptx.save(file_path)
        
        logger.info(f"PowerPointファイル生成完了: {file_path}")
        return file_path
    
    except Exception as e:
        logger.error(f"PowerPointファイル生成エラー: {str(e)}")
        logger.debug(f"詳細なエラー: {traceback.format_exc()}")
        
        # エラー時にはデモファイルを返す
        try:
            # デモ用の簡易ファイルを作成
            pptx = PPTXPresentation()
            slide = pptx.slides.add_slide(pptx.slide_layouts[0])
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = "デモプレゼンテーション"
            
            # サブタイトル
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # サブタイトル
                    shape.text = "エラーが発生したため、デモファイルを生成しました"
                    break
            
            # エラースライド
            slide = pptx.slides.add_slide(pptx.slide_layouts[1])
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = "エラー情報"
            
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"エラーが発生しました: {str(e)}"
            
            # 一時ファイルに保存
            file_name = f"demo_presentation_{presentation.id}.pptx"
            file_path = os.path.join(TEMP_DIR, file_name)
            pptx.save(file_path)
            
            logger.info(f"デモファイル生成完了: {file_path}")
            return file_path
        except Exception as demo_error:
            logger.error(f"デモファイル生成エラー: {str(demo_error)}")
            raise ValueError(f"PowerPointファイルの生成中にエラーが発生し、デモファイルの作成にも失敗しました: {str(e)}")


def get_theme_settings(theme: str) -> Dict[str, Any]:
    """テーマに基づいた設定を取得する"""
    themes = {
        "modern": {
            "primary_color": "0EA5E9",
            "secondary_color": "7DD3FC",
            "text_color": "374151",
            "title_font": "Arial",
            "content_font": "Arial",
        },
        "business": {
            "primary_color": "1E40AF",
            "secondary_color": "3B82F6",
            "text_color": "1F2937",
            "title_font": "Calibri",
            "content_font": "Calibri",
        },
        "creative": {
            "primary_color": "8B5CF6",
            "secondary_color": "C4B5FD",
            "text_color": "4B5563",
            "title_font": "Segoe UI",
            "content_font": "Segoe UI",
        },
        "minimal": {
            "primary_color": "64748B",
            "secondary_color": "94A3B8",
            "text_color": "334155",
            "title_font": "Helvetica",
            "content_font": "Helvetica",
        },
    }
    
    logger.debug(f"テーマ {theme} の設定を取得")
    return themes.get(theme, themes["modern"])
